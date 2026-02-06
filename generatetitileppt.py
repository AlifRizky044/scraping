from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Data
data = [
    ("2025-09-06 19:18:59", "Heri Julianto", "089648192255"),
    ("2025-09-06 19:19:14", "Wira jaya", "08999423206"),
    ("2025-09-06 19:19:17", "Arman Arroisi", "089764015101"),
    ("2025-09-06 19:19:46", "Wahyu Putra Mulyadi", "089654966044"),
    ("2025-09-06 19:20:07", "Yumanita", "0895343266159"),
    ("2025-09-06 19:20:14", "Putra dewa ramadhan Ibrahim", "089695646830"),
    ("2025-09-06 19:20:32", "Muhammad arasy", "0897 9722 195"),
    ("2025-09-06 19:20:35", "Ilham", "0895382511441"),
    ("2025-09-06 19:20:39", "Takiishi", "08992690902"),
    ("2025-09-06 19:20:41", "Fandi Ahmad Fauzan", "0895706035566"),
    ("2025-09-06 19:21:15", "Istara Putra Artha Wicaksana", "089668213899"),
    ("2025-09-06 19:21:17", "M. Fariq Nafis Zakaria", "089620559559"),
    ("2025-09-06 19:21:22", "Rizky aditya", "089677758689"),
    ("2025-09-06 19:21:27", "HILMI AULIA LUKMAN", "0895332843460"),
    ("2025-09-06 19:21:59", "Cahyo Heru Wicaksono", "0895393124667"),
    ("2025-09-06 19:22:02", "DICKY HAMZAH", "089531093474"),
    ("2025-09-06 19:22:03", "Renaldi wicaksana", "0895377751976"),
    ("2025-09-06 19:22:05", "RHOBET ARNOLD HALOMOAN NABABAN", "089514449155"),
    ("2025-09-06 19:22:08", "Niko jelang ramadhan isnaen pratama", "08998458040"),
]

# Create presentation
prs = Presentation()

for date, name, phone in data:
    slide_layout = prs.slide_layouts[5]  # Blank slide
    slide = prs.slides.add_slide(slide_layout)

    # Table (1 header row + 1 data row)
    rows, cols = 2, 3
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(2)).table

    # Column widths
    for i in range(3):
        table.columns[i].width = Inches(3)

    # Header
    headers = ["Date & Time", "Name", "Phone Number"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Data row
    row_data = [date, name, phone]
    for col, val in enumerate(row_data):
        table.cell(1, col).text = val
        for paragraph in table.cell(1, col).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0, 0, 0)

# Save
prs.save("contacts_slides.pptx")
print("contacts_slides.pptx generated!")


# cara run program scrapping
# source venv/bin/activate  
# python generatetitileppt.py  
